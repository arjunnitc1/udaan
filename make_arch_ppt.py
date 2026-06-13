from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette ──────────────────────────────────────────────────────────────
INK      = RGBColor(0x1B, 0x21, 0x43)
INK2     = RGBColor(0x2A, 0x31, 0x60)
INK3     = RGBColor(0x3A, 0x42, 0x70)
MARIGOLD = RGBColor(0xF2, 0xA0, 0x07)
MARIGOLD2= RGBColor(0xD8, 0x8B, 0x00)
ROSE     = RGBColor(0xC8, 0x50, 0x6E)
PAPER    = RGBColor(0xFF, 0xFD, 0xF7)
SAND     = RGBColor(0xF6, 0xEF, 0xE2)
LEAF     = RGBColor(0x2E, 0x7D, 0x5B)
INDIGO   = RGBColor(0x43, 0x38, 0xCA)
INDIGO2  = RGBColor(0x63, 0x66, 0xF1)
SKY      = RGBColor(0x0E, 0xA5, 0xE9)
GOLD     = RGBColor(0xD9, 0x77, 0x06)
LINE     = RGBColor(0xE8, 0xDF, 0xCC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GHOST    = RGBColor(0xC0, 0xC8, 0xFF)
DIM      = RGBColor(0x80, 0x88, 0xAA)
RED      = RGBColor(0xDC, 0x26, 0x26)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Primitives ─────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=13, bold=False, color=INK,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Palatino Linotype" if (bold and size >= 22) else font
    return tb

def para(tf, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT,
         italic=False, before=Pt(3)):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = before
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color; r.font.name = "Calibri"
    return p

def kite(slide, x, y, sz=Inches(0.5), col=MARIGOLD):
    s = slide.shapes.add_shape(1, x, y, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = col
    s.line.fill.background(); s.rotation = 45

def dark_bg(slide): rect(slide, 0, 0, W, H, fill=INK)
def paper_bg(slide): rect(slide, 0, 0, W, H, fill=PAPER)
def eyebrow(slide, text, x, y, col=MARIGOLD2):
    txt(slide, text.upper(), x, y, Inches(12), Inches(0.35),
        size=8, bold=True, color=col)

def label_box(slide, text, x, y, w, h, bg, fg, size=10, bold=True):
    rect(slide, x, y, w, h, fill=bg)
    txt(slide, text, x + Inches(0.12), y + Inches(0.08),
        w - Inches(0.24), h - Inches(0.12), size=size, bold=bold,
        color=fg, align=PP_ALIGN.CENTER)

def arrow_right(slide, x, y, length=Inches(0.5), col=MARIGOLD):
    line = slide.shapes.add_shape(9, x, y, length, Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = col
    line.line.fill.background()

def connector(slide, x1, y1, x2, y2, col=DIM, lw=Pt(1.5)):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = col; c.line.width = lw

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
rect(s, W - Inches(4), 0, Inches(4), H, fill=INK2)
rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=MARIGOLD)
kite(s, Inches(10.0), Inches(0.5), Inches(1.2), MARIGOLD)
kite(s, Inches(11.6), Inches(1.8), Inches(0.6), ROSE)
kite(s, Inches(10.8), Inches(5.9), Inches(0.8), RGBColor(0x43,0x38,0xCA))

txt(s, "UDAAN  उड़ान", Inches(0.6), Inches(0.8), Inches(8.5), Inches(0.5),
    size=11, bold=True, color=MARIGOLD)
txt(s, "Technical Architecture", Inches(0.6), Inches(1.4), Inches(8.5), Inches(1.5),
    size=46, bold=True, color=WHITE)
txt(s, "Deep Dive", Inches(0.6), Inches(2.9), Inches(8.5), Inches(0.9),
    size=36, bold=True, color=MARIGOLD, italic=True)

txt(s, "How we built a production-grade AI coaching platform\nfor 197 million Indian women — from first principles.",
    Inches(0.6), Inches(4.1), Inches(8.5), Inches(1.2),
    size=15, color=GHOST)

pillars = [
    (MARIGOLD,  "Anti-hallucination\nby design"),
    (ROSE,      "Structured AI\noutput, not prose"),
    (INDIGO2,   "Zero vendor lock-in\ndatabase"),
    (LEAF,      "Voice-native\nmultilingual"),
    (SKY,       "Privacy-first\n(DPDP ready)"),
]
for i, (col, label) in enumerate(pillars):
    cx = Inches(0.6) + i * Inches(1.72)
    rect(s, cx, Inches(5.6), Inches(1.6), Inches(1.55), fill=INK2)
    rect(s, cx, Inches(5.6), Inches(1.6), Inches(0.05), fill=col)
    txt(s, label, cx + Inches(0.1), Inches(5.72), Inches(1.4), Inches(1.2),
        size=9, bold=True, color=col, align=PP_ALIGN.CENTER)

txt(s, "AI for Good Hackathon  ·  Reuben College, Oxford  ·  June 2026",
    Inches(0.6), Inches(7.15), Inches(8), Inches(0.3), size=9, color=DIM)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture Philosophy
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
paper_bg(s)
rect(s, 0, 0, Inches(0.2), H, fill=MARIGOLD)

eyebrow(s, "Architecture Philosophy — Three Principles", Inches(0.5), Inches(0.4))
txt(s, "Every decision has a reason. Here's ours.", Inches(0.5), Inches(0.82),
    Inches(12), Inches(0.75), size=28, bold=True, color=INK)

principles = [
    (MARIGOLD, "01", "AI does structured reasoning,\nnot free conversation.",
     "The model returns a fully typed JSON object (the Kit) — not chat text. Every field is validated. "
     "The UI renders data, not prose. This means the output is verifiable, saveable, and testable. "
     "Most 'AI apps' at hackathons return a string. We return a type-safe business plan."),
    (ROSE,     "02", "Ground truth lives in\ncode, not in model weights.",
     "Government schemes, pricing rules, and coaching ethics are encoded in files and prompts — not "
     "trusted to the model's memory. schemes.json is the source of truth. The prompt injects it at "
     "runtime. Claude cannot hallucinate a scheme that isn't in the file."),
    (INDIGO2,  "03", "The database is infrastructure,\nnot a dependency.",
     "We use Drizzle ORM with a single DATABASE_URL environment variable. Switching from Supabase "
     "to our own self-deployed PostgreSQL requires changing one line. No application code changes. "
     "This is intentional: we built for data sovereignty from day one."),
]

for i, (col, num, headline, body) in enumerate(principles):
    cy = Inches(1.85) + i * Inches(1.72)
    rect(s, Inches(0.45), cy, Inches(12.4), Inches(1.6), fill=WHITE,
         line=LINE, lw=Pt(1))
    rect(s, Inches(0.45), cy, Inches(0.06), Inches(1.6), fill=col)
    txt(s, num, Inches(0.65), cy + Inches(0.15), Inches(0.7), Inches(0.5),
        size=22, bold=True, color=col)
    txt(s, headline, Inches(0.65), cy + Inches(0.62), Inches(3.6), Inches(0.85),
        size=13, bold=True, color=INK)
    txt(s, body, Inches(4.5), cy + Inches(0.18), Inches(8.2), Inches(1.28),
        size=11, color=INK3)
    rect(s, Inches(4.3), cy + Inches(0.3), Inches(0.03), Inches(1.0), fill=LINE)

txt(s, "These three principles are why our architecture looks different from other teams'.",
    Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.35),
    size=10, italic=True, color=ROSE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Full System Architecture Diagram
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
kite(s, Inches(12.5), Inches(0.3), Inches(0.7), MARIGOLD)

eyebrow(s, "System Architecture — Full Stack", Inches(0.5), Inches(0.38))
txt(s, "Three layers. One data flow. Zero black boxes.", Inches(0.5), Inches(0.78),
    Inches(10), Inches(0.65), size=24, bold=True, color=WHITE)

# Layer 1: Browser
rect(s, Inches(0.4), Inches(1.55), Inches(12.5), Inches(1.55), fill=INK2)
rect(s, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.05), fill=SKY)
txt(s, "LAYER 1 — BROWSER  (Next.js 14 App Router · Client Components)",
    Inches(0.6), Inches(1.6), Inches(8), Inches(0.38), size=10, bold=True, color=SKY)

browser_boxes = [
    (SKY,     "Landing Page\n/"),
    (SKY,     "Auth\n/auth"),
    (SKY,     "Dashboard\n/dashboard"),
    (MARIGOLD,"Coach Chat\n/coach"),
    (SKY,     "Community\n/community"),
    (SKY,     "Goals\n/goals"),
    (SKY,     "Piggy Bank\n/piggy-bank"),
]
for i, (col, label) in enumerate(browser_boxes):
    cx = Inches(0.55) + i * Inches(1.77)
    rect(s, cx, Inches(2.04), Inches(1.68), Inches(0.92), fill=RGBColor(0x0E,0x20,0x40))
    rect(s, cx, Inches(2.04), Inches(1.68), Inches(0.04), fill=col)
    txt(s, label, cx + Inches(0.08), Inches(2.1), Inches(1.52), Inches(0.8),
        size=8.5, bold=True, color=col, align=PP_ALIGN.CENTER)

# Web APIs strip
rect(s, Inches(0.4), Inches(3.18), Inches(12.5), Inches(0.42), fill=RGBColor(0x05,0x10,0x28))
for i, api in enumerate(["Web Speech API (voice input · 6 langs)",
                          "SpeechSynthesis TTS (voice output · zero cost)",
                          "localStorage (auth · goals · piggy bank · notifications)"]):
    cx = Inches(0.6) + i * Inches(4.1)
    txt(s, "⚡ " + api, cx, Inches(3.24), Inches(3.9), Inches(0.3),
        size=8, color=GHOST)

# Arrow down
txt(s, "↓  HTTPS  ↓", Inches(0.4), Inches(3.65), Inches(12.5), Inches(0.32),
    size=9, bold=True, color=DIM, align=PP_ALIGN.CENTER)

# Layer 2: API
rect(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(1.48), fill=INK2)
rect(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.05), fill=MARIGOLD)
txt(s, "LAYER 2 — API LAYER  (Next.js Route Handlers · Server-side only)",
    Inches(0.6), Inches(4.05), Inches(8), Inches(0.38), size=10, bold=True, color=MARIGOLD)

api_boxes = [
    (MARIGOLD, "/api/coach", "POST\nZod validate → Rate limit → Claude → Parse JSON → Persist Kit"),
    (SKY,      "/api/session", "GET/POST\nFunnel analytics: session_start → interview_answer → kit_generated → whatsapp_copied"),
    (ROSE,     "Upstash Redis", "Sliding window\n10 req / 60s per IP · Serverless-safe REST API"),
    (INDIGO2,  "Anthropic SDK", "claude-opus-4-8\nServer-side only · API key never in browser"),
]
for i, (col, title, body) in enumerate(api_boxes):
    cx = Inches(0.55) + i * Inches(3.1)
    rect(s, cx, Inches(4.5), Inches(2.95), Inches(0.85), fill=RGBColor(0x0A,0x14,0x2E))
    rect(s, cx, Inches(4.5), Inches(2.95), Inches(0.04), fill=col)
    txt(s, title, cx + Inches(0.1), Inches(4.56), Inches(2.75), Inches(0.32),
        size=10, bold=True, color=col)
    txt(s, body, cx + Inches(0.1), Inches(4.9), Inches(2.75), Inches(0.38),
        size=7.5, color=GHOST)

# Arrow down
txt(s, "↓  PostgreSQL wire protocol  ↓", Inches(0.4), Inches(5.52), Inches(12.5), Inches(0.3),
    size=9, bold=True, color=DIM, align=PP_ALIGN.CENTER)

# Layer 3: Data
rect(s, Inches(0.4), Inches(5.86), Inches(12.5), Inches(1.3), fill=RGBColor(0x08,0x28,0x18))
rect(s, Inches(0.4), Inches(5.86), Inches(12.5), Inches(0.05), fill=LEAF)
txt(s, "LAYER 3 — DATA LAYER  (PostgreSQL via Drizzle ORM · migrating to self-hosted)",
    Inches(0.6), Inches(5.9), Inches(9), Inches(0.35), size=10, bold=True, color=LEAF)

db_boxes = [
    (LEAF, "sessions table", "UUID PK · lang · is_demo · created_at\nNo PII stored — ever."),
    (LEAF, "events table",   "session_id FK · type · payload JSONB · created_at\nFunnel analytics, cascade delete."),
    (LEAF, "kits table",     "session_id FK · kit_json JSONB · created_at\nFull Kit object per completed session."),
    (MARIGOLD2, "schemes.json", "4 verified schemes · injected at runtime\nSingle source of truth — not model weights."),
]
for i, (col, title, body) in enumerate(db_boxes):
    cx = Inches(0.55) + i * Inches(3.1)
    rect(s, cx, Inches(6.3), Inches(2.95), Inches(0.78), fill=RGBColor(0x05,0x18,0x0E))
    txt(s, title, cx + Inches(0.1), Inches(6.34), Inches(2.75), Inches(0.3),
        size=10, bold=True, color=col)
    txt(s, body, cx + Inches(0.1), Inches(6.64), Inches(2.75), Inches(0.38),
        size=7.5, color=GHOST)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Request Flow (sequence diagram)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
paper_bg(s)
rect(s, 0, 0, Inches(0.2), H, fill=INDIGO)

eyebrow(s, "Request Flow — What Happens on Every Coach Message", Inches(0.5), Inches(0.4), col=INDIGO)
txt(s, "7 steps, 3 layers, one typed Kit. Nothing is skipped.", Inches(0.5), Inches(0.82),
    Inches(12), Inches(0.65), size=26, bold=True, color=INK)

steps = [
    (INDIGO2, "1",  "User sends message",
     "Browser validates non-empty input. Attaches `x-session-id` header (client UUID — no PII). "
     "Sends conversation history as a JSON array of `{role, content}` pairs to POST /api/coach."),
    (SKY,     "2",  "Zod schema validation",
     "Server validates every message: role must be 'user' | 'assistant', content 1–2000 chars, "
     "max 20 messages in history. Malformed requests return 400 before touching Claude."),
    (ROSE,    "3",  "Rate limit check (Upstash Redis)",
     "Sliding window: 10 requests per 60 seconds per IP address. Returns 429 with Retry-After "
     "headers if exceeded. Skipped gracefully if Upstash env vars absent (local dev)."),
    (MARIGOLD,"4",  "System prompt injection",
     "buildSystemPrompt() constructs the full prompt at request time — injecting all 4 verified "
     "schemes from schemes.json. Prompt includes conversation rules, language detection, JSON "
     "output format, and the labour-as-cost pricing constraint."),
    (MARIGOLD,"5",  "Claude API call",
     "Anthropic SDK sends messages + system prompt. Model: claude-opus-4-8. max_tokens: 1500. "
     "No streaming — waits for complete JSON before responding to browser."),
    (LEAF,    "6",  "JSON parse + defensive cleaning",
     "Response text has code fences stripped (```json ... ``` → raw JSON). Primary: JSON.parse(). "
     "Fallback: regex extract first `{...}` block. If both fail, returns 500 with user-friendly error."),
    (LEAF,    "7",  "Persist Kit + return",
     "If reply.type === 'kit', inserts kit_json (full typed Kit object) into kits table linked to "
     "session_id. Kit persist failure is non-fatal — logs error, still returns Kit to user. "
     "Always returns {reply, raw} to browser."),
]

for i, (col, num, title, body) in enumerate(steps):
    cy = Inches(1.75) + i * Inches(0.79)
    rect(s, Inches(0.45), cy, Inches(0.48), Inches(0.7), fill=col)
    txt(s, num, Inches(0.45), cy + Inches(0.1), Inches(0.48), Inches(0.5),
        size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(1.05), cy + Inches(0.04), Inches(3.5), Inches(0.38),
        size=11, bold=True, color=INK)
    txt(s, body, Inches(4.7), cy + Inches(0.04), Inches(8.0), Inches(0.65),
        size=9.5, color=INK3)
    rect(s, Inches(4.5), cy + Inches(0.1), Inches(0.03), Inches(0.5), fill=LINE)

txt(s, "Each step has a reason. Each failure mode has a fallback. Nothing is fire-and-forget except DB writes (non-fatal by design).",
    Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
    size=9.5, italic=True, color=ROSE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AI Pipeline: 5 Layers
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
kite(s, Inches(12.3), Inches(6.2), Inches(0.9), MARIGOLD)

eyebrow(s, "AI Pipeline — Why This Is Not a Chatbot Wrapper", Inches(0.5), Inches(0.38))
txt(s, "Five engineering layers on top of the model call.", Inches(0.5), Inches(0.78),
    Inches(10), Inches(0.65), size=26, bold=True, color=WHITE)

layers = [
    (MARIGOLD, "Structured Conversation Protocol",
     "The model is not free-chatting. Each response returns a JSON object with a `type` field: "
     "'question' (with a `stage` tag: life|motivation|skills|business|detail) or 'kit'. "
     "The UI transitions state based on type. This is a state machine, not a chatroom.",
     "Code: app/coach/page.tsx → sendMsg() → reply.type switch"),

    (SKY,      "Anti-Hallucination Scheme Grounding",
     "schemes.json is the ONLY source of scheme data. buildSystemPrompt() injects it as plain text "
     "at runtime with the instruction: 'Use ONLY these schemes, never invent others.' "
     "Claude has no way to use a scheme that isn't in the file. This is structural, not instructional.",
     "Code: lib/prompt.ts → buildSystemPrompt() · data/schemes.json"),

    (ROSE,     "Labour-As-Cost Pricing Constraint",
     "A hard rule enforced in every prompt call: 'Always count her own labour as a cost when pricing. "
     "Women underprice because they don't value their time — you must.' This is an economic "
     "intervention, not a feature. It runs on every Kit generation, with no opt-out.",
     "Code: lib/prompt.ts → system prompt, line 18"),

    (LEAF,     "Structured JSON Output (typed Kit)",
     "The model is instructed to return ONLY valid JSON matching the Kit schema (15+ fields, "
     "bilingual _en/_hi variants). The response is parsed, JSON.parsed, and typed as Kit. "
     "If it fails parsing, there's a regex fallback. If that fails, a 500 error — never a raw string.",
     "Code: app/api/coach/route.ts → steps 4–5 · lib/types.ts → Kit type"),

    (INDIGO2,  "Multilingual Auto-Detection + TTS",
     "User input is scanned for Unicode ranges: Devanagari (hi), Tamil (ta), Telugu (te), "
     "Kannada (kn), Bengali (bn). Detected language code is passed to SpeechSynthesisUtterance. "
     "No external TTS API, no cost, no latency. The browser speaks the response.",
     "Code: app/coach/page.tsx → detectLangCode() · speak()"),
]

for i, (col, title, body, code) in enumerate(layers):
    cy = Inches(1.65) + i * Inches(1.12)
    rect(s, Inches(0.45), cy, Inches(0.06), Inches(0.92), fill=col)
    txt(s, title, Inches(0.65), cy + Inches(0.06), Inches(4.0), Inches(0.42),
        size=12, bold=True, color=col)
    txt(s, body, Inches(0.65), cy + Inches(0.48), Inches(7.6), Inches(0.5),
        size=9.5, color=GHOST)
    rect(s, Inches(8.5), cy + Inches(0.12), Inches(4.5), Inches(0.72),
         fill=RGBColor(0x0A,0x14,0x2E), line=col, lw=Pt(0.5))
    txt(s, "📎  " + code, Inches(8.62), cy + Inches(0.24), Inches(4.26), Inches(0.5),
        size=8.5, color=col, italic=True)

txt(s, "A 'chatbot wrapper' returns a string. Udaan returns a type-safe, validated, renderable data structure.",
    Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.35),
    size=10, italic=True, color=MARIGOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The Anti-Hallucination Engine
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
paper_bg(s)
rect(s, 0, 0, Inches(0.2), H, fill=MARIGOLD)

eyebrow(s, "Anti-Hallucination Layer — Schemes Grounding Deep Dive", Inches(0.5), Inches(0.4))
txt(s, "The model cannot invent a scheme that doesn't exist in our code.",
    Inches(0.5), Inches(0.82), Inches(12), Inches(0.72), size=26, bold=True, color=INK)

# Left: the flow
flow_items = [
    (MARIGOLD2, "schemes.json",
     "4 verified government schemes. Each entry has:\nid · name · what_en/hi · fits[] · step_en/hi · official_url · last_verified"),
    (INK,       "buildSystemPrompt()",
     "At every API call, iterates schemes and formats them as\nplain text: name | what | Hindi | first step | fits | URL"),
    (INDIGO2,   "System Prompt",
     "Injected text: 'VERIFIED GOVERNMENT SCHEMES — use ONLY these,\nnever invent others.' Hard constraint, not a suggestion."),
    (ROSE,      "Claude Response",
     "Model selects 2–3 schemes from the injected list that genuinely\nfit the user's profile. Cannot cite anything outside the list."),
    (LEAF,      "Kit.schemes[]",
     "Typed array of scheme objects rendered in the UI with\nname, explanation, and first concrete step."),
]
for i, (col, title, body) in enumerate(flow_items):
    cy = Inches(1.75) + i * Inches(1.02)
    rect(s, Inches(0.45), cy, Inches(5.8), Inches(0.88), fill=WHITE, line=LINE, lw=Pt(1))
    rect(s, Inches(0.45), cy, Inches(0.06), Inches(0.88), fill=col)
    txt(s, title, Inches(0.65), cy + Inches(0.08), Inches(3.0), Inches(0.35),
        size=11, bold=True, color=col)
    txt(s, body, Inches(0.65), cy + Inches(0.46), Inches(5.1), Inches(0.35),
        size=9, color=INK3)
    if i < 4:
        txt(s, "↓", Inches(2.8), cy + Inches(0.92), Inches(0.5), Inches(0.12),
            size=10, color=DIM, align=PP_ALIGN.CENTER)

# Right: scheme cards
rect(s, Inches(7.0), Inches(1.68), Inches(5.9), Inches(5.5),
     fill=SAND, line=LINE, lw=Pt(1))
txt(s, "What's actually in schemes.json", Inches(7.15), Inches(1.74),
    Inches(5.6), Inches(0.38), size=10, bold=True, color=MARIGOLD2)

schemes_data = [
    ("Udyam Registration",  "Free MSME registration in ~15 min.\nOpens loans + subsidies.",  "any small business"),
    ("PM Mudra Loan",       "Collateral-free loan up to ₹50,000\n(Shishu tier) for startups.", "needs capital"),
    ("SHG / NRLM linkage",  "Join local women's group: savings,\nlow-interest credit, network.", "rural · wants community"),
    ("PM Vishwakarma",      "For artisans/tailors: training,\ntoolkit grant, credit.",         "tailoring · crafts"),
]
for i, (name, what, fits) in enumerate(schemes_data):
    cy = Inches(2.22) + i * Inches(1.2)
    rect(s, Inches(7.15), cy, Inches(5.6), Inches(1.08),
         fill=WHITE, line=LINE, lw=Pt(1))
    txt(s, name, Inches(7.28), cy + Inches(0.08), Inches(5.2), Inches(0.35),
        size=10, bold=True, color=INK)
    txt(s, what, Inches(7.28), cy + Inches(0.44), Inches(5.2), Inches(0.35),
        size=8.5, color=INK3)
    rect(s, Inches(7.28), cy + Inches(0.8), Inches(5.2), Inches(0.22), fill=SAND)
    txt(s, "Fits: " + fits, Inches(7.35), cy + Inches(0.82), Inches(5.0), Inches(0.18),
        size=7.5, bold=True, color=MARIGOLD2)

txt(s, "Each scheme has official_url and last_verified date. We update the file — Claude adapts automatically.",
    Inches(0.45), Inches(7.08), Inches(12.3), Inches(0.38),
    size=9.5, italic=True, color=ROSE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — The Database Decision
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
kite(s, Inches(12.0), Inches(0.4), Inches(0.8), MARIGOLD)

eyebrow(s, "The Database Decision — Supabase → Self-Deployed PostgreSQL", Inches(0.5), Inches(0.38))
txt(s, "Why we're migrating. Why it's easy. Why it matters.", Inches(0.5), Inches(0.78),
    Inches(10), Inches(0.65), size=26, bold=True, color=WHITE)

# Before / After columns
rect(s, Inches(0.45), Inches(1.58), Inches(5.9), Inches(5.12), fill=INK2)
rect(s, Inches(0.45), Inches(1.58), Inches(5.9), Inches(0.06), fill=RED)
txt(s, "BEFORE — Supabase (managed)", Inches(0.6), Inches(1.65),
    Inches(5.6), Inches(0.38), size=11, bold=True, color=RED)

before = [
    "✗  US-based servers — potential DPDP Act issue",
    "✗  Supabase session-mode pooler required\n   (max:1, prepare:false workarounds in code)",
    "✗  Free tier: 500MB storage, 2 projects",
    "✗  Vendor lock-in — if Supabase changes\n   pricing, we're exposed",
    "✗  Cannot add pgvector for future\n   RAG scheme matching",
    "✗  Connection limits on serverless fns\n   forced our architecture",
]
for i, item in enumerate(before):
    cy = Inches(2.12) + i * Inches(0.72)
    txt(s, item, Inches(0.62), cy, Inches(5.6), Inches(0.62),
        size=9.5, color=GHOST)

rect(s, Inches(6.95), Inches(1.58), Inches(6.0), Inches(5.12), fill=RGBColor(0x08,0x28,0x18))
rect(s, Inches(6.95), Inches(1.58), Inches(6.0), Inches(0.06), fill=LEAF)
txt(s, "AFTER — Self-Deployed PostgreSQL", Inches(7.1), Inches(1.65),
    Inches(5.6), Inches(0.38), size=11, bold=True, color=LEAF)

after = [
    "✓  India-hosted server — DPDP Act compliant",
    "✓  Direct connection, no pooler workarounds\n   (remove max:1, prepare:false)",
    "✓  Unlimited storage on our own instance",
    "✓  Zero vendor lock-in — we own the data\n   and the infrastructure",
    "✓  Add pgvector extension → enables full\n   RAG scheme matching pipeline",
    "✓  PgBouncer for connection pooling —\n   production-grade, our control",
]
for i, item in enumerate(after):
    cy = Inches(2.12) + i * Inches(0.72)
    txt(s, item, Inches(7.1), cy, Inches(5.7), Inches(0.62),
        size=9.5, color=GHOST)

# Migration complexity
rect(s, Inches(0.45), Inches(6.78), Inches(12.5), Inches(0.6), fill=MARIGOLD)
txt(s, "Migration complexity: Change ONE environment variable (DATABASE_URL). "
       "Zero application code changes. Drizzle ORM abstracts the connection completely. "
       "Run drizzle-kit push on new host. Done.",
    Inches(0.6), Inches(6.83), Inches(12.2), Inches(0.5),
    size=10.5, bold=True, color=INK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Voice & Language Pipeline
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
paper_bg(s)
rect(s, 0, 0, Inches(0.2), H, fill=SKY)

eyebrow(s, "Voice & Language Pipeline — Zero External Services", Inches(0.5), Inches(0.4), col=SKY)
txt(s, "6 Indian languages. No AWS Polly. No Google TTS. No API cost.",
    Inches(0.5), Inches(0.82), Inches(12), Inches(0.68), size=25, bold=True, color=INK)

# Pipeline flow
pipeline = [
    (SKY,     "Voice Input\n(Web Speech API)",
     "Browser's SpeechRecognition API.\nLanguage: hi-IN or en-IN\nbased on UI toggle.",
     "webkitSpeechRecognition\ninterimResults: true\nNo server call"),
    (INDIGO2, "Language Detection\n(Unicode Ranges)",
     "detectLangCode() scans input\nfor Devanagari, Tamil, Telugu,\nKannada, Bengali characters.",
     "Regex: /[\\u0900-\\u097F]/\nfor Hindi etc.\nZero latency"),
    (MARIGOLD,"Claude API Call\n(Multilingual prompt)",
     "Prompt instructs model to respond\nin detected language. Model returns\nbilingual Kit (en + hi always).",
     "System prompt:\n'Respond in whatever\nlanguage user writes in'"),
    (LEAF,    "TTS Output\n(SpeechSynthesis API)",
     "speak() creates\nSpeechSynthesisUtterance\nwith detected lang code.",
     "SpeechSynthesisUtterance\n.lang = detected code\nVoice matched by browser"),
]

for i, (col, title, body, code) in enumerate(pipeline):
    cx = Inches(0.45) + i * Inches(3.15)
    rect(s, cx, Inches(1.75), Inches(3.02), Inches(3.8), fill=WHITE, line=LINE, lw=Pt(1))
    rect(s, cx, Inches(1.75), Inches(3.02), Inches(0.05), fill=col)
    txt(s, title, cx + Inches(0.15), Inches(1.85), Inches(2.72), Inches(0.65),
        size=11, bold=True, color=col)
    txt(s, body, cx + Inches(0.15), Inches(2.56), Inches(2.72), Inches(1.5),
        size=10, color=INK3)
    rect(s, cx + Inches(0.1), Inches(4.05), Inches(2.82), Inches(1.3),
         fill=SAND, line=LINE, lw=Pt(0.5))
    txt(s, code, cx + Inches(0.18), Inches(4.12), Inches(2.65), Inches(1.15),
        size=8.5, color=INK, italic=True)
    if i < 3:
        txt(s, "→", Inches(0.45) + i * Inches(3.15) + Inches(3.08),
            Inches(3.3), Inches(0.2), Inches(0.4),
            size=16, bold=True, color=col, align=PP_ALIGN.CENTER)

# Languages table
rect(s, Inches(0.45), Inches(5.75), Inches(12.5), Inches(1.42), fill=INK)
txt(s, "Languages supported", Inches(0.6), Inches(5.82),
    Inches(3), Inches(0.35), size=10, bold=True, color=MARIGOLD)

langs = [
    ("Hindi", "hi-IN", "हिन्दी", "Devanagari U+0900–U+097F"),
    ("Tamil", "ta-IN", "தமிழ்", "Tamil U+0B80–U+0BFF"),
    ("Telugu", "te-IN", "తెలుగు", "Telugu U+0C00–U+0C7F"),
    ("Kannada", "kn-IN", "ಕನ್ನಡ", "Kannada U+0C80–U+0CFF"),
    ("Bengali", "bn-IN", "বাংলা", "Bengali U+0980–U+09FF"),
    ("English", "en-IN", "English", "Default fallback"),
]
for i, (name, code, native, urange) in enumerate(langs):
    cx = Inches(0.6) + i * Inches(2.05)
    txt(s, native, cx, Inches(6.22), Inches(1.95), Inches(0.38),
        size=14, bold=True, color=MARIGOLD, align=PP_ALIGN.CENTER)
    txt(s, name + " · " + code, cx, Inches(6.6), Inches(1.95), Inches(0.25),
        size=8, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, urange, cx, Inches(6.85), Inches(1.95), Inches(0.25),
        size=7, color=DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Data Architecture & Privacy
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
kite(s, Inches(12.2), Inches(0.4), Inches(0.8), MARIGOLD)

eyebrow(s, "Data Architecture — Privacy By Design  ·  DPDP Act 2023 Ready", Inches(0.5), Inches(0.38))
txt(s, "What we store. What we don't. Why every choice was deliberate.",
    Inches(0.5), Inches(0.78), Inches(12), Inches(0.65), size=24, bold=True, color=WHITE)

# What we store
rect(s, Inches(0.45), Inches(1.58), Inches(5.85), Inches(5.4), fill=INK2)
rect(s, Inches(0.45), Inches(1.58), Inches(5.85), Inches(0.06), fill=LEAF)
txt(s, "WHAT WE STORE  (minimal, purposeful)", Inches(0.6), Inches(1.65),
    Inches(5.5), Inches(0.38), size=10, bold=True, color=LEAF)

stores = [
    ("sessions", "UUID · language · is_demo · timestamp",
     "No name, no phone, no location. UUID generated client-side."),
    ("events",   "session_id · event_type · payload · timestamp",
     "Funnel analytics only: session_start, interview_answer,\nkit_generated, whatsapp_copied."),
    ("kits",     "session_id · kit_json · timestamp",
     "Full Kit for product analytics. No user identity attached.\nCascade deletes with session."),
    ("schemes.json","Static file · version-controlled",
     "Government schemes. No user data. Deployed with code.\nUpdated manually, re-verified each cycle."),
]
for i, (table, schema, purpose) in enumerate(stores):
    cy = Inches(2.15) + i * Inches(1.18)
    rect(s, Inches(0.62), cy, Inches(5.5), Inches(1.0), fill=RGBColor(0x0A,0x14,0x2E))
    txt(s, table, Inches(0.75), cy + Inches(0.08), Inches(1.8), Inches(0.3),
        size=11, bold=True, color=LEAF)
    txt(s, schema, Inches(0.75), cy + Inches(0.38), Inches(5.2), Inches(0.28),
        size=8.5, color=MARIGOLD, italic=True)
    txt(s, purpose, Inches(0.75), cy + Inches(0.65), Inches(5.2), Inches(0.3),
        size=8.5, color=GHOST)

# What we don't store
rect(s, Inches(6.95), Inches(1.58), Inches(5.9), Inches(5.4), fill=RGBColor(0x22,0x0A,0x0A))
rect(s, Inches(6.95), Inches(1.58), Inches(5.9), Inches(0.06), fill=RED)
txt(s, "WHAT WE NEVER STORE", Inches(7.1), Inches(1.65),
    Inches(5.5), Inches(0.38), size=10, bold=True, color=RED)

no_stores = [
    ("❌ Phone number", "Auth is localStorage only. No server ever sees it."),
    ("❌ Real name", "All sessions use client-generated UUID."),
    ("❌ Location", "Not collected at any point in the flow."),
    ("❌ Financial data", "Income projections live in Kit JSON, linked to no identity."),
    ("❌ Conversation text", "Messages are never persisted. Only Kit output is stored."),
    ("❌ Device identifiers", "Rate limiting uses IP only, not stored beyond Redis TTL."),
    ("❌ Photos or media", "Text and voice only. No uploads."),
]
for i, (what, why) in enumerate(no_stores):
    cy = Inches(2.15) + i * Inches(0.69)
    txt(s, what, Inches(7.12), cy, Inches(2.5), Inches(0.32),
        size=10, bold=True, color=RED)
    txt(s, why, Inches(9.72), cy + Inches(0.02), Inches(2.8), Inches(0.32),
        size=9, color=GHOST)

# DPDP footer
rect(s, Inches(0.45), Inches(7.06), Inches(12.5), Inches(0.38), fill=RGBColor(0x1A,0x2A,0x0A))
txt(s, "India DPDP Act 2023: No 'personal data' (as defined in the Act) is collected or processed. "
       "Self-hosted DB migration completes the data sovereignty picture.",
    Inches(0.6), Inches(7.1), Inches(12.2), Inches(0.3),
    size=9, color=LEAF, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Security Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
paper_bg(s)
rect(s, 0, 0, Inches(0.2), H, fill=ROSE)

eyebrow(s, "Security Architecture — API Key Never Leaves the Server", Inches(0.5), Inches(0.4), col=ROSE)
txt(s, "Designed to be safe when serving vulnerable users at scale.",
    Inches(0.5), Inches(0.82), Inches(12), Inches(0.65), size=26, bold=True, color=INK)

sec_items = [
    (ROSE,    "🔑  API Key Isolation",
     "ANTHROPIC_API_KEY is a server-side environment variable. It is never in any client bundle, "
     "never in any response header, never logged. The browser makes requests to /api/coach — "
     "a Next.js Route Handler that runs server-side only. The key is structurally inaccessible to the browser.",
     "Verify: Network tab → /api/coach response never contains key"),
    (MARIGOLD,"🛡️  Input Validation (Zod)",
     "Every API request is validated before any processing. CoachRequestSchema enforces: "
     "role ∈ {user, assistant}, content length 1–2000 chars, max 20 messages. "
     "SessionEventSchema enforces event types. Malformed requests return 400 — Claude is never called.",
     "Code: lib/validations.ts · used in both route handlers"),
    (INDIGO2, "⏱️  Rate Limiting (Upstash Redis)",
     "Sliding window algorithm: 10 requests per 60 seconds per IP. Returns 429 with standard "
     "RateLimit headers (Limit, Remaining, Reset). Graceful degradation: if Redis is unavailable, "
     "the limiter is skipped rather than breaking the app (fail-open by design for MVP).",
     "Code: app/api/coach/route.ts · Upstash sliding window"),
    (LEAF,    "🗄️  Database Security",
     "No PII in database — UUID-keyed sessions only. Session IDs are client-generated "
     "(Math.random() UUID) so we cannot map a session back to a phone number even if compromised. "
     "Kit persist is non-fatal: DB errors are caught, logged, and the response still returns to user.",
     "Code: lib/db/schema.ts · no PII columns"),
    (SKY,     "🔒  Auth Architecture",
     "Current: localStorage-based session (MVP). No server-side session = no session hijacking "
     "surface. Phone number and OTP never leave the device. Moving to: httpOnly cookie session "
     "with CSRF protection as next step before production scale.",
     "Planned: next-auth or custom httpOnly cookie + server session"),
]

for i, (col, title, body, note) in enumerate(sec_items):
    cy = Inches(1.65) + i * Inches(1.05)
    rect(s, Inches(0.45), cy, Inches(12.45), Inches(0.96),
         fill=WHITE, line=LINE, lw=Pt(1))
    rect(s, Inches(0.45), cy, Inches(0.06), Inches(0.96), fill=col)
    txt(s, title, Inches(0.65), cy + Inches(0.06), Inches(3.8), Inches(0.42),
        size=11, bold=True, color=col)
    txt(s, body, Inches(0.65), cy + Inches(0.5), Inches(7.6), Inches(0.42),
        size=9, color=INK3)
    rect(s, Inches(8.5), cy + Inches(0.1), Inches(4.2), Inches(0.76),
         fill=SAND, line=LINE, lw=Pt(0.5))
    txt(s, note, Inches(8.62), cy + Inches(0.18), Inches(4.0), Inches(0.6),
        size=8, color=INK3, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Performance & Scale
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
kite(s, Inches(12.1), Inches(0.3), Inches(0.8), MARIGOLD)

eyebrow(s, "Performance & Scalability — Designed for Serverless India", Inches(0.5), Inches(0.38))
txt(s, "Every architectural choice has a scale implication.", Inches(0.5), Inches(0.78),
    Inches(10), Inches(0.65), size=26, bold=True, color=WHITE)

perf_items = [
    (MARIGOLD, "Lazy DB singleton",
     "Database client created on first request, not at module import. This prevents "
     "connection pool exhaustion in serverless environments (Vercel, AWS Lambda) "
     "where many function instances start cold simultaneously.",
     "Result: Zero connection errors on cold start · No 'too many clients' errors"),
    (SKY, "max:1 connection per serverless fn",
     "Each serverless invocation gets exactly 1 Postgres connection. "
     "When we move to self-hosted PgBouncer, this becomes configurable — "
     "we can raise it per function based on load patterns.",
     "Result: Predictable connection count · No pool exhaustion at scale"),
    (ROSE, "Non-fatal DB writes",
     "Kit persistence is wrapped in try/catch. If the database is slow or down, "
     "the user still gets their Kit. Analytics are best-effort — the product "
     "never degrades because of observability failures.",
     "Result: User-facing latency unaffected by DB slowdowns"),
    (LEAF, "No streaming (by design for MVP)",
     "The API waits for the complete Claude response before returning. This simplifies "
     "error handling (parse the full JSON once) and avoids partial Kit renders. "
     "Streaming can be added later for conversational feel without changing the Kit logic.",
     "Result: Simple, reliable JSON parse · No partial state bugs"),
    (INDIGO2, "Static pages where possible",
     "Community heroes, landing page, auth — all static (no DB calls). "
     "Only /api/coach and /api/session are dynamic. Next.js build shows 8 static routes, "
     "2 dynamic API routes. CDN-cacheable pages for 80% of traffic.",
     "Result: Instant loads for most pages · DB only hit for coaching + analytics"),
]

for i, (col, title, body, result) in enumerate(perf_items):
    cy = Inches(1.65) + i * Inches(1.08)
    rect(s, Inches(0.45), cy, Inches(0.06), Inches(0.9), fill=col)
    txt(s, title, Inches(0.65), cy + Inches(0.06), Inches(6.5), Inches(0.38),
        size=11, bold=True, color=col)
    txt(s, body, Inches(0.65), cy + Inches(0.46), Inches(6.5), Inches(0.5),
        size=9.5, color=GHOST)
    rect(s, Inches(7.4), cy + Inches(0.1), Inches(5.5), Inches(0.76),
         fill=RGBColor(0x0A,0x14,0x2E), line=col, lw=Pt(0.5))
    txt(s, result, Inches(7.52), cy + Inches(0.2), Inches(5.26), Inches(0.52),
        size=9, color=col, italic=False, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Future Architecture (Roadmap)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
paper_bg(s)
rect(s, 0, 0, Inches(0.2), H, fill=INDIGO)

eyebrow(s, "Architecture Roadmap — What The Foundation Enables", Inches(0.5), Inches(0.4), col=INDIGO)
txt(s, "The current design isn't the destination. It's the launchpad.",
    Inches(0.5), Inches(0.82), Inches(12), Inches(0.65), size=26, bold=True, color=INK)

# Timeline: Now / 3mo / 6mo / 12mo
cols_data = [
    (LEAF, "Now — Built", [
        "✓ Next.js 14 App Router",
        "✓ Claude conversational coach",
        "✓ Anti-hallucination schemes",
        "✓ Structured Kit JSON",
        "✓ 6 Indian languages + TTS",
        "✓ Funnel analytics",
        "✓ Goals + Piggy Bank",
        "✓ Community heroes",
    ]),
    (SKY, "3 months", [
        "→ Self-hosted PostgreSQL",
        "→ pgvector + RAG scheme matching",
        "→ Deterministic pricing engine",
        "→ httpOnly cookie auth",
        "→ WhatsApp Bot (Meta Cloud API)",
        "→ Push notifications (FCM)",
        "→ Gemini Flash fallback",
    ]),
    (MARIGOLD, "6 months", [
        "→ More Indian languages (mr, gu, pa)",
        "→ NGO/SHG partner dashboard",
        "→ First-sale tracking metric",
        "→ Order management module",
        "→ Basic inventory tracker",
        "→ UPI payment links in Kit",
    ]),
    (ROSE, "12 months", [
        "→ Fine-tuned model on Indian\n   women entrepreneur corpus",
        "→ Marketplace integration\n   (Meesho, Flipkart Samarth)",
        "→ Multilingual voice-only mode\n   for low-literacy users",
        "→ Offline PWA for rural India",
    ]),
]

for i, (col, period, items) in enumerate(cols_data):
    cx = Inches(0.45) + i * Inches(3.18)
    rect(s, cx, Inches(1.72), Inches(3.05), Inches(5.22), fill=WHITE, line=LINE, lw=Pt(1))
    rect(s, cx, Inches(1.72), Inches(3.05), Inches(0.06), fill=col)
    txt(s, period, cx + Inches(0.15), Inches(1.82), Inches(2.75), Inches(0.42),
        size=12, bold=True, color=col)
    for j, item in enumerate(items):
        cy = Inches(2.35) + j * Inches(0.56)
        txt(s, item, cx + Inches(0.15), cy, Inches(2.75), Inches(0.52),
            size=9, color=INK3 if col != LEAF else RGBColor(0x08,0x50,0x28))

# Key unlock: pgvector
rect(s, Inches(0.45), Inches(7.05), Inches(12.5), Inches(0.4), fill=INK)
txt(s, "Key unlock: Self-hosted PostgreSQL + pgvector → real RAG scheme matching → "
       "scores eligibility in code, not LLM → verifiable, auditable, faster.",
    Inches(0.6), Inches(7.1), Inches(12.2), Inches(0.3),
    size=10, bold=True, color=MARIGOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Closing
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
rect(s, 0, 0, W, Inches(0.06), fill=MARIGOLD)
rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=MARIGOLD)
kite(s, Inches(5.7), Inches(0.5), Inches(2.0), MARIGOLD)
kite(s, Inches(7.2), Inches(1.8), Inches(1.0), ROSE)
kite(s, Inches(6.4), Inches(2.8), Inches(0.6), INDIGO2)

txt(s, "Udaan  उड़ान", Inches(0.8), Inches(3.0), Inches(12), Inches(1.2),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Anti-hallucination  ·  Structured AI output  ·  Data sovereignty  ·  Voice-native  ·  Privacy-first",
    Inches(0.8), Inches(4.3), Inches(12), Inches(0.5),
    size=13, color=MARIGOLD, align=PP_ALIGN.CENTER, italic=True)
txt(s, "github.com/arjunnitc1/udaan  ·  branch: shreya/scheme-engine",
    Inches(0.8), Inches(5.3), Inches(12), Inches(0.5),
    size=14, bold=True, color=GHOST, align=PP_ALIGN.CENTER)
txt(s, "AI for Good Hackathon  ·  Reuben College, Oxford  ·  June 2026",
    Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
    size=10, color=DIM, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "/Users/shreyasinha/Desktop/udaan/docs/udaan-architecture-deep-dive.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
